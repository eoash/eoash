"""컨텍스트 엔지니어링 교육 슬라이드 1장 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

BG = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
ACCENT = RGBColor(0xE8, 0xFF, 0x47)  # lime
CARD_BG = RGBColor(0x1A, 0x1A, 0x1A)
RED_DIM = RGBColor(0xFF, 0x66, 0x66)
GREEN_DIM = RGBColor(0x66, 0xFF, 0x99)

# Background
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = BG

def add_box(slide, left, top, width, height, text, font_size=18,
            color=WHITE, bold=False, align=PP_ALIGN.LEFT, fill=None, anchor=None):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    if fill:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Malgun Gothic"
    return txBox

def add_rounded_card(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle as card background"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

# Helpers
def inch(v): return int(v * 914400)
def cm(v): return int(v * 360000)

W = inch(13.333)
H = inch(7.5)
MARGIN = inch(0.8)

# ── Title ──
add_box(slide, MARGIN, inch(0.4), inch(10), inch(0.6),
        "컨텍스트 엔지니어링 = 파일 설계", 36, ACCENT, bold=True)

add_box(slide, MARGIN, inch(1.0), inch(10), inch(0.4),
        "채팅 AI는 블랙박스, Claude Code는 파일 시스템", 18, GRAY)

# ── Two comparison cards ──
card_w = inch(5.4)
card_h = inch(2.2)
card_y = inch(1.7)
gap = inch(0.6)
left_x = MARGIN
right_x = MARGIN + card_w + gap

# Left card: 채팅 AI
add_rounded_card(slide, left_x, card_y, card_w, card_h, RGBColor(0x1A, 0x0A, 0x0A))
add_box(slide, left_x + inch(0.3), card_y + inch(0.2), card_w, inch(0.4),
        "채팅 AI", 22, RED_DIM, bold=True)
add_box(slide, left_x + inch(0.3), card_y + inch(0.7), card_w - inch(0.6), inch(0.35),
        "컨텍스트 = 블랙박스", 18, WHITE)
add_box(slide, left_x + inch(0.3), card_y + inch(1.1), card_w - inch(0.6), inch(0.35),
        "대화 기록이 어떻게 저장되는지 모름", 14, GRAY)
add_box(slide, left_x + inch(0.3), card_y + inch(1.5), card_w - inch(0.6), inch(0.35),
        "설정할 수 없고, 확인할 수 없음", 14, GRAY)

# Right card: Claude Code
add_rounded_card(slide, right_x, card_y, card_w, card_h, RGBColor(0x0A, 0x1A, 0x0A))
add_box(slide, right_x + inch(0.3), card_y + inch(0.2), card_w, inch(0.4),
        "Claude Code", 22, GREEN_DIM, bold=True)
add_box(slide, right_x + inch(0.3), card_y + inch(0.7), card_w - inch(0.6), inch(0.35),
        "컨텍스트 = 파일", 18, WHITE)
add_box(slide, right_x + inch(0.3), card_y + inch(1.1), card_w - inch(0.6), inch(0.35),
        "직접 열어보고, 수정할 수 있음", 14, GRAY)
add_box(slide, right_x + inch(0.3), card_y + inch(1.5), card_w - inch(0.6), inch(0.35),
        "구조를 직접 설계할 수 있음", 14, GRAY)

# ── Divider line ──
div_y = inch(4.15)
from pptx.enum.shapes import MSO_SHAPE
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Emu(div_y), Emu(W - MARGIN*2), Emu(inch(0.01)))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
line.line.fill.background()

# ── Section title ──
add_box(slide, MARGIN, inch(4.3), inch(8), inch(0.45),
        "컨텍스트 설계 3원칙", 24, WHITE, bold=True)

# ── Three principle cards ──
p_card_w = inch(3.6)
p_card_h = inch(2.2)
p_card_y = inch(5.0)
p_gap = inch(0.4)

principles = [
    ("간결성", "CLAUDE.md", "96줄 — 매 세션 읽히는 지도\n짧을수록 좋다 (토큰 = 비용)"),
    ("계층화", "identity.md 등", "상세 정보는 별도 파일에\nCLAUDE.md엔 링크만 기록"),
    ("피드백 루프", "ANTI_PATTERNS.md", "AI 실수를 기록 → 오답노트\n다음 세션부터 반복 안 함"),
]

for i, (title, file_name, desc) in enumerate(principles):
    px = MARGIN + i * (p_card_w + p_gap)

    add_rounded_card(slide, px, p_card_y, p_card_w, p_card_h, CARD_BG)

    # Number circle effect
    add_box(slide, px + inch(0.3), p_card_y + inch(0.15), inch(0.4), inch(0.4),
            str(i+1), 20, BG, bold=True)
    # Fake circle bg
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(px + inch(0.25)), Emu(p_card_y + inch(0.18)),
        Emu(inch(0.35)), Emu(inch(0.35))
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    # Number on top of circle
    add_box(slide, px + inch(0.25), p_card_y + inch(0.18), inch(0.35), inch(0.35),
            str(i+1), 16, BG, bold=True, align=PP_ALIGN.CENTER)

    # Principle title
    add_box(slide, px + inch(0.7), p_card_y + inch(0.2), p_card_w - inch(1), inch(0.35),
            title, 20, ACCENT, bold=True)

    # File name
    add_box(slide, px + inch(0.3), p_card_y + inch(0.7), p_card_w - inch(0.6), inch(0.3),
            file_name, 13, RGBColor(0xAA, 0xAA, 0xAA))

    # Description (multiline)
    desc_box = add_box(slide, px + inch(0.3), p_card_y + inch(1.05), p_card_w - inch(0.6), inch(1.0),
            "", 14, GRAY)
    tf = desc_box.text_frame
    tf.word_wrap = True
    for j, line_text in enumerate(desc.split("\n")):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(14)
        run.font.color.rgb = GRAY
        run.font.name = "Malgun Gothic"
        p.space_after = Pt(4)

# Save
out_path = "ai-native-camp/training/slide_context_engineering.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
