#!/usr/bin/env python3
"""AI Native 교육 슬라이드 생성 — training/ai_native_intro.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colors ───────────────────────────────────────────────────
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
INDIGO = RGBColor(0x5B, 0x6C, 0xF9)
CARD   = RGBColor(0xF5, 0xF5, 0xFA)
IND_L  = RGBColor(0xEB, 0xED, 0xFF)
GRAY   = RGBColor(0x9A, 0x9A, 0xA8)
WARN   = RGBColor(0xFF, 0x5C, 0x35)
DK_BG  = RGBColor(0x1A, 0x1B, 0x2E)

FONT = "Pretendard"


def make_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Helpers ──────────────────────────────────────────────

    def blank(bg=WHITE):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        f = s.background.fill
        f.solid()
        f.fore_color.rgb = bg
        return s

    def box(s, x, y, w, h, fill, line_color=None):
        sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line_color:
            sh.line.color.rgb = line_color
        else:
            sh.line.fill.background()
        return sh

    def tb(s, text, x, y, w, h, size, bold=False, color=DARK, align=PP_ALIGN.LEFT):
        t = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = t.text_frame
        tf.word_wrap = True
        lines = text.split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FONT
        return t

    def bar(s):
        box(s, 0, 0, 13.33, 0.07, INDIGO)

    def title(s, text, color=DARK):
        tb(s, text, 0.6, 0.32, 12.13, 0.75, 27, bold=True, color=color)

    def card(s, x, y, w, h, icon, ttl, body, tc=DARK, bc=GRAY, bg=CARD):
        box(s, x, y, w, h, bg)
        tb(s, icon, x+0.22, y+0.2,  0.65, 0.55, 26)
        tb(s, ttl,  x+0.22, y+0.85, w-0.44, 0.55, 15, bold=True, color=tc)
        tb(s, body, x+0.22, y+1.5,  w-0.44, h-1.75, 12, color=bc)

    # ── Slide 1: 오프닝 ──────────────────────────────────────
    s1 = blank()
    box(s1, 0, 0, 0.18, 7.5, INDIGO)
    tb(s1, "우리는 AI를",        0.75, 1.5,  11.5, 1.3, 58, bold=True, color=DARK)
    tb(s1, "잘못 쓰고 있다",     0.75, 2.85, 11.5, 1.3, 58, bold=True, color=INDIGO)
    tb(s1, "ChatGPT, Gemini, Claude — 다들 쓰고 있다.\n근데 왜 업무가 별로 안 바뀐 것 같지?",
       0.75, 4.45, 11.0, 1.0, 17, color=GRAY)

    # ── Slide 2: 경각심 ──────────────────────────────────────
    s2 = blank(); bar(s2)
    title(s2, "지금 밖에서 무슨 일이 벌어지고 있나")
    CW, CY, CH = 3.84, 1.38, 5.1
    CX = [0.6, 4.77, 8.94]
    card(s2, CX[0], CY, CW, CH, "💬",
         '"토큰에 변화가 없던데?"',
         '설에 쉬는 회사 대표에게\n다른 회사 대표가 날린 한마디\n\n— Claude Code League 中')
    card(s2, CX[1], CY, CW, CH, "🏢",
         '인티그레이션 · 157명',
         '✓ 150명 (95%) 전환 완료\n✓ 마케터·인사·정산팀이 상위권\n✓ 정산팀이 자동화 툴 직접 개발\n✓ BO가 사내 시스템 직접 배포')
    card(s2, CX[2], CY, CW, CH, "⚠️",
         '임계점',
         '임계점을 넘은 팀원과\n못 넘은 팀원은\n\n"대화가 되지 않습니다"',
         tc=WARN, bc=DARK)

    # ── Slide 3: 4단계 진화 ───────────────────────────────────
    s3x = blank(); bar(s3x)
    title(s3x, "AI 활용은 4단계로 진화하고 있다")

    # 4 columns
    EW = 2.88          # card width
    EH = 4.65          # card height
    EY = 1.38
    EXS = [0.6 + i * (EW + 0.25) for i in range(4)]

    stages = [
        ("1", "Prompt\nEngineering",
         "채팅창에 질문하기",
         '"이 코드 에러 원인이 뭐야?"',
         "단발성 질의 · 창 닫으면 끝"),
        ("2", "Context\nEngineering",
         "CLAUDE.md로 맞춤 세팅",
         "프로젝트 규칙·스타일을\nAI가 매번 기억",
         "한 번 설정 → 계속 적용"),
        ("3", "Compound\nEngineering",
         "팀 공유 + 외부 도구 결합",
         "Slack·Notion·DB 스크립트를\nAI가 직접 실행·조합",
         "개인 → 조직 레벨 자동화"),
        ("4", "Harness\nEngineering",
         "자율 에이전틱 루프",
         "코드 작성 → 테스트 → 에러 수정\n루프를 AI가 스스로 반복",
         "24시간 풀 자동화 · CI/CD 결합"),
    ]

    for i, (num, name, sub, desc, bottom) in enumerate(stages):
        is_first = (i == 0)
        cbg   = RGBColor(0xFF, 0xF0, 0xEB) if is_first else CARD
        cline = WARN if is_first else None
        c = box(s3x, EXS[i], EY, EW, EH, cbg, line_color=cline)

        # stage number circle
        num_color = WARN if is_first else INDIGO
        tb(s3x, num, EXS[i]+0.15, EY+0.15, 0.5, 0.5, 22, bold=True, color=num_color)

        # stage name
        name_color = WARN if is_first else INDIGO
        tb(s3x, name, EXS[i]+0.15, EY+0.65, EW-0.3, 0.85, 16, bold=True, color=name_color)

        # subtitle
        tb(s3x, sub, EXS[i]+0.15, EY+1.55, EW-0.3, 0.5, 11, bold=True, color=DARK)

        # description
        tb(s3x, desc, EXS[i]+0.15, EY+2.15, EW-0.3, 1.2, 11, color=GRAY)

        # bottom label
        tb(s3x, bottom, EXS[i]+0.15, EY+3.55, EW-0.3, 0.9, 10, color=GRAY)

        # arrow between cards
        if i < 3:
            tb(s3x, "→", EXS[i]+EW+0.02, EY+EH/2-0.25, 0.21, 0.5, 14,
               bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # "YOU ARE HERE" marker on stage 1
    box(s3x, EXS[0], EY+EH+0.12, EW, 0.55, WARN)
    tb(s3x, "대부분이 여기에 있다", EXS[0], EY+EH+0.18, EW, 0.42,
       12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # bottom message
    box(s3x, 0.6, 6.55, 12.13, 0.7, IND_L)
    tb(s3x, "채팅 AI에 머물면 임계점을 넘을 수 없다",
       0.6, 6.65, 12.13, 0.5, 16, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    # ── Slide 4 (was 3): 기억 못한다 ──────────────────────────
    s3 = blank(); bar(s3)
    title(s3, "AI도, 슬랙도 — 기억하지 못한다")
    card(s3, 0.6, 1.38, 5.93, 4.6, "🤖",
         '채팅 AI',
         '"분명히 설명했는데 왜 또 물어봐?"\n\n→ 대화창 닫으면 초기화\n→ 컨텍스트가 길어질수록\n   앞 내용을 잊는다')
    card(s3, 6.8, 1.38, 5.93, 4.6, "💬",
         '슬랙',
         '"3개월 전 대화, 찾을 수 있나요?"\n\n→ 빠른 소통의 대가는 휘발\n→ 일주일 전 업무 요청,\n   지금 추적 가능한가요?')
    box(s3, 0.6, 6.18, 12.13, 0.87, IND_L)
    tb(s3, "빠름의 대가가 '휘발'이다",
       0.6, 6.32, 12.13, 0.6, 20, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    # ── Slide 4: Claude Code 비교 ────────────────────────────
    s4 = blank(); bar(s4)
    title(s4, "터미널을 켜면 AI의 속이 보인다")
    box(s4, 0.6, 1.38, 5.55, 4.95, CARD)
    tb(s4, "기존 채팅 방식", 0.9, 1.6, 5.0, 0.55, 15, bold=True, color=GRAY)
    for i, t in enumerate(["대화 → 사라짐", "프롬프트가 중요", "매번 설명 반복", "창 닫으면 초기화"]):
        tb(s4, f"✕  {t}", 0.9, 2.38 + i*0.9, 5.0, 0.6, 14, color=GRAY)
    tb(s4, "→", 6.35, 3.5, 0.63, 0.75, 30, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    box(s4, 7.18, 1.38, 5.55, 4.95, IND_L)
    tb(s4, "Claude Code", 7.48, 1.6, 5.0, 0.55, 15, bold=True, color=INDIGO)
    for i, t in enumerate(["대화 → 파일로 저장", "컨텍스트가 중요", "한 번 설정 → 계속 기억", "작업 히스토리가 쌓인다"]):
        tb(s4, f"✓  {t}", 7.48, 2.38 + i*0.9, 5.0, 0.6, 14, color=INDIGO)

    # ── Slide 5: Context Engineering ─────────────────────────
    s5 = blank(); bar(s5)
    title(s5, "프롬프트 잘 쓰는 시대는 끝났다")
    box(s5, 0.6, 1.4, 5.5, 1.95, CARD)
    tb(s5, "Prompt Engineering",  0.88, 1.58, 5.0, 0.55, 14, bold=True, color=GRAY)
    tb(s5, '"질문을 잘 써야 해"', 0.88, 2.18, 5.0, 0.55, 13, color=GRAY)
    tb(s5, "→", 6.35, 2.1, 0.63, 0.8, 28, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    box(s5, 7.23, 1.4, 5.5, 1.95, IND_L)
    tb(s5, "Context Engineering", 7.5, 1.58, 5.0, 0.55, 14, bold=True, color=INDIGO)
    tb(s5, '"정보를 잘 쌓아야 해"', 7.5, 2.18, 5.0, 0.55, 13, color=INDIGO)
    box(s5, 0.6, 3.68, 12.13, 3.05, INDIGO)
    tb(s5, "잘 정리된 Notion 페이지",     1.0, 3.9,  11.0, 0.85, 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s5, "=",                           1.0, 4.72, 11.0, 0.6,  20,            color=WHITE, align=PP_ALIGN.CENTER)
    tb(s5, "AI에게 주는 좋은 컨텍스트", 1.0, 5.28, 11.0, 0.85, 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 6: 복리 + 임계점 ───────────────────────────────
    s6 = blank(); bar(s6)
    title(s6, "매일 쌓이면 복리가 된다")
    SW, SY, SH = 2.82, 1.38, 4.18
    steps = [
        ("Day 1",  "CLAUDE.md 만들기"),
        ("Day 7",  "AI가 내 업무\n스타일을 파악"),
        ("Day 30", '"이거 저번에\n어떻게 했지?"\nAI가 기억한다'),
        ("Day 90", "내 업무 히스토리가\n완성된다"),
    ]
    SX = [0.6 + i * (SW + 0.3) for i in range(4)]
    for i, (day, desc) in enumerate(steps):
        bg = IND_L if i == 3 else CARD
        box(s6, SX[i], SY, SW, SH, bg)
        tb(s6, day,  SX[i]+0.2, SY+0.2,    SW-0.4, 0.6,      18, bold=True, color=INDIGO)
        tb(s6, desc, SX[i]+0.2, SY+0.95,   SW-0.4, SH-1.2,   13, color=DARK)
        if i < 3:
            tb(s6, "→", SX[i]+SW+0.03, SY+SH/2-0.3, 0.24, 0.55, 13,
               bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    box(s6, 0.6, 5.75, 12.13, 1.1, CARD)
    tb(s6, "임계점을 넘으면, 일하는 방식이 달라진다",
       0.6, 5.93, 12.13, 0.75, 18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ── Slide 7: Phase 1 ─────────────────────────────────────
    s7 = blank(); bar(s7)
    title(s7, "Phase 1 — 나만의 업무 기록부터")
    CX7 = [0.6, 4.77, 8.94]
    card(s7, CX7[0], 1.38, 3.84, 4.6, "📄",
         "내 CLAUDE.md",
         "= 내 업무 설명서\nAI가 나를 기억하는 파일\n나만의 컨텍스트")
    card(s7, CX7[1], 1.38, 3.84, 4.6, "📌",
         "매일 commit",
         "= 내 업무 일지\n누가 시키지 않아도\n기록이 쌓인다")
    card(s7, CX7[2], 1.38, 3.84, 4.6, "⚡",
         "/wrap",
         "= 자동 정리\n퇴근 전 한 번만 치면\nAI가 알아서 정리")
    box(s7, 0.6, 6.18, 12.13, 0.87, IND_L)
    tb(s7, "내일 아침, Claude Code를 켜는 것부터 시작입니다",
       0.6, 6.32, 12.13, 0.6, 17, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    # ── Slide 8: Phase 2 ─────────────────────────────────────
    s8 = blank(); bar(s8)
    title(s8, "Phase 2 — 팀이 함께 쌓으면")
    tb(s8, "공유 지능", 0.6, 1.22, 12.13, 0.9, 34, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    pairs = [
        ('"지은님이 저번 주에 뭐 하셨죠?"', "AI가 알고 있다"),
        ('"이 업무 누가 담당하고 있죠?"',   "AI가 알고 있다"),
        ('"비슷한 거 예전에 해봤나요?"',    "AI가 찾아준다"),
    ]
    for i, (q, a) in enumerate(pairs):
        cy = 2.32 + i * 1.48
        box(s8, 0.6, cy, 8.5, 1.25, CARD)
        tb(s8, q, 0.85, cy+0.3, 8.0, 0.65, 14, color=DARK)
        box(s8, 9.33, cy, 3.4, 1.25, IND_L)
        tb(s8, a, 9.33, cy+0.3, 3.4, 0.65, 14, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    tb(s8, "슬랙 대화가 아니라, git 히스토리가 팀의 기억이 됩니다",
       0.6, 6.85, 12.13, 0.5, 13, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 9: 행동 지시 (dark) ────────────────────────────
    s9 = blank(bg=DK_BG)
    tb(s9, "오늘 저녁부터",
       1.2, 0.9,  11.0, 0.75, 22, color=GRAY,  align=PP_ALIGN.CENTER)
    tb(s9, "딱 이것만",
       1.2, 1.55, 11.0, 1.1,  46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    acts = [
        ("①", "Claude Code 설치"),
        ("②", "내일 아침, PC 켜면 Claude Code도 켠다"),
        ("③", "퇴근 전  /wrap  →  commit  →  push  →  PC 끈다"),
    ]
    for i, (num, text) in enumerate(acts):
        cy = 3.05 + i * 1.22
        tb(s9, num,  1.5,  cy+0.08, 0.85, 0.9, 30, bold=True, color=INDIGO)
        tb(s9, text, 2.55, cy+0.15, 9.5,  0.75, 20, color=WHITE)

    # ── Save ─────────────────────────────────────────────────
    out = os.path.join(os.path.dirname(__file__), "ai_native_intro.pptx")
    prs.save(out)
    print(f"저장 완료: {out}")
    print(f"슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    make_pptx()
