#!/usr/bin/env python3
"""AI Native 교육 슬라이드 v2 — 14장 (training/ai_native_intro_v2.pptx)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Colors ───────────────────────────────────────────────────
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
INDIGO = RGBColor(0x5B, 0x6C, 0xF9)
CARD   = RGBColor(0xF5, 0xF5, 0xFA)
IND_L  = RGBColor(0xEB, 0xED, 0xFF)
GRAY   = RGBColor(0x9A, 0x9A, 0xA8)
WARN   = RGBColor(0xFF, 0x5C, 0x35)
DK_BG  = RGBColor(0x1A, 0x1B, 0x2E)
WARN_L = RGBColor(0xFF, 0xF0, 0xEB)
GREEN  = RGBColor(0x2D, 0xB5, 0x6D)

FONT = "Pretendard"
SW = 13.33  # slide width inches
SH = 7.5    # slide height inches


def make_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    # ── Helpers ──────────────────────────────────────────────

    def blank(bg=WHITE):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        f = s.background.fill
        f.solid()
        f.fore_color.rgb = bg
        return s

    def box(s, x, y, w, h, fill, line_color=None):
        sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line_color:
            sh.line.color.rgb = line_color
            sh.line.width = Pt(2)
        else:
            sh.line.fill.background()
        return sh

    def tb(s, text, x, y, w, h, size, bold=False, color=DARK, align=PP_ALIGN.LEFT, spacing=None):
        t = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = t.text_frame
        tf.word_wrap = True
        lines = text.split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if spacing:
                p.space_after = Pt(spacing)
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FONT
        return t

    def bar(s):
        box(s, 0, 0, SW, 0.07, INDIGO)

    def title(s, text, color=DARK):
        tb(s, text, 0.6, 0.32, 12.13, 0.75, 27, bold=True, color=color)

    def card(s, x, y, w, h, icon, ttl, body, tc=DARK, bc=GRAY, bg=CARD):
        box(s, x, y, w, h, bg)
        tb(s, icon, x+0.22, y+0.2,  0.65, 0.55, 26)
        tb(s, ttl,  x+0.22, y+0.85, w-0.44, 0.55, 15, bold=True, color=tc)
        tb(s, body, x+0.22, y+1.5,  w-0.44, h-1.75, 12, color=bc)

    def bottom_bar(s, text, bg=IND_L, color=INDIGO, y=6.55, h=0.7, size=16):
        box(s, 0.6, y, 12.13, h, bg)
        tb(s, text, 0.6, y+0.1, 12.13, h-0.2, size, bold=True, color=color, align=PP_ALIGN.CENTER)

    # ── Slide 1: 오프닝 ──────────────────────────────────────
    s = blank()
    box(s, 0, 0, 0.18, SH, INDIGO)
    tb(s, "우리는 AI를",        0.75, 1.5,  11.5, 1.3, 58, bold=True, color=DARK)
    tb(s, "잘못 쓰고 있다",     0.75, 2.85, 11.5, 1.3, 58, bold=True, color=INDIGO)
    tb(s, "ChatGPT, Gemini, Claude — 다들 쓰고 있다.\n근데 왜 업무가 별로 안 바뀐 것 같지?",
       0.75, 4.45, 11.0, 1.0, 17, color=GRAY)

    # ── Slide 2a: METR 벤치마크 ──────────────────────────────
    s = blank(); bar(s)
    title(s, "AI가 혼자 일할 수 있는 시간")
    tb(s, '새 모델이 나올 때마다 "어떤 모델이 1등인가"만 보고 있지 않나요?\n그건 우리에게 직접 와닿지 않습니다.\n이 그래프는 다릅니다 — AI가 사람의 일을 얼마나 대체할 수 있는지를 보여줍니다.',
       0.6, 1.15, 12.13, 1.0, 12, color=GRAY)

    # Bar chart - horizontal bars
    data = [
        ("2019  GPT-2",           0.04,    "2초",      0.3),
        ("2023  GPT-4",           3.5,     "3분",      0.5),
        ("2024  Claude 3.5",      11,      "11분",     0.8),
        ("2024  o1",              38,      "38분",     1.5),
        ("2025  Claude 3.7",      60,      "1시간",    2.2),
        ("2025  Opus 4.5",        320,     "5시간",    4.5),
        ("2026  Opus 4.6",        870,     "14.5시간", 9.5),
    ]
    chart_x = 0.6
    chart_y = 2.35
    chart_w = 11.5
    row_h = 0.52
    max_val = 870

    for i, (label, val, display, bar_w) in enumerate(data):
        ry = chart_y + i * (row_h + 0.12)
        # label
        tb(s, label, chart_x, ry, 3.0, row_h, 11, color=DARK)
        # bar
        bx = chart_x + 3.2
        is_last = (i == len(data) - 1)
        bar_color = INDIGO if is_last else RGBColor(0xC5, 0xCB, 0xFC)
        box(s, bx, ry+0.05, bar_w, row_h-0.1, bar_color)
        # value text
        tb(s, display, bx + bar_w + 0.15, ry, 1.5, row_h, 11, bold=is_last,
           color=INDIGO if is_last else GRAY)

    tb(s, "7개월마다 2배로 성장  — METR Horizon Benchmark",
       0.6, 6.15, 12.13, 0.35, 11, color=GRAY, align=PP_ALIGN.RIGHT)
    bottom_bar(s, "조만간 일주일이 되고, 한 달이 되고, 1년이 됩니다.")

    # ── Slide 2b: 해커톤 ─────────────────────────────────────
    s = blank(); bar(s)
    title(s, "지금 밖에서 무슨 일이 벌어지고 있나")
    tb(s, "Claude Code 해커톤 수상자 명단에 개발자가 거의 없습니다",
       0.6, 1.15, 12.13, 0.5, 14, color=GRAY)

    # try to add image
    img_path = r"c:\Users\ash\OneDrive\바탕 화면\claude code 해커톤 수상자 명단에 개발자가 한 명도 없습니다..jpg"
    try:
        s.shapes.add_picture(img_path, Inches(0.6), Inches(1.75), Inches(5.8), Inches(3.8))
    except Exception:
        box(s, 0.6, 1.75, 5.8, 3.8, CARD)
        tb(s, "[해커톤 캡처 이미지]", 0.6, 3.2, 5.8, 1.0, 14, color=GRAY, align=PP_ALIGN.CENTER)

    # right side - winner list
    rx = 6.8
    winners = [
        ("🥇 금상",     "CrossBeam",    "캘리포니아 상해 전문 변호사"),
        ("🥈 은상",     "Elisa",        "12살 딸을 둔 아버지"),
        ("🥉 동상",     "postvisit.ai", "벨기에 심장내과 과장"),
        ("🎨 Creative", "Conductr",     "에스토니아 뮤지션"),
        ("🧠 Keep",     "TARA",         "우간다 인프라 엔지니어"),
    ]
    for i, (medal, name, job) in enumerate(winners):
        ry = 1.85 + i * 0.82
        box(s, rx, ry, 5.93, 0.7, CARD)
        tb(s, medal, rx+0.15, ry+0.12, 1.5, 0.46, 11, bold=True, color=INDIGO)
        tb(s, name,  rx+1.7,  ry+0.05, 1.8, 0.3, 11, bold=True, color=DARK)
        tb(s, job,   rx+1.7,  ry+0.35, 4.0, 0.3, 10, color=WARN)

    bottom_bar(s, '"어떻게 만드는가"보다 "무엇이 필요한가"가 핵심 경쟁력이 되고 있다')

    # ── Slide 2c: 토큰 전쟁 ──────────────────────────────────
    s = blank(); bar(s)
    title(s, "한국에서도 이미 이렇게 일하고 있다")
    tb(s, "의료 플랫폼 '메디스트림' vs AX 기업 '렛서' — Claude Code 토큰 전쟁",
       0.6, 1.15, 12.13, 0.5, 14, color=GRAY)

    # try to add image
    img_token = r"c:\Users\ash\OneDrive\바탕 화면\토큰전쟁1.jpg"
    try:
        s.shapes.add_picture(img_token, Inches(0.6), Inches(1.85), Inches(5.5), Inches(4.2))
    except Exception:
        box(s, 0.6, 1.85, 5.5, 4.2, CARD)
        tb(s, "[토큰전쟁 캡처 이미지]", 0.6, 3.5, 5.5, 1.0, 14, color=GRAY, align=PP_ALIGN.CENTER)

    # right side - key facts
    facts = [
        ("95%", "157명 중 150명이\nClaude Code로 전환"),
        ("비개발자", "토큰 최상위권 =\n마케터·인사·정산·디자이너"),
        ("직역 파괴", "기획이 개발하고\n디자이너가 프론트를 배포"),
        ("임계점", "넘은 팀원과 못 넘은 팀원은\n대화가 되지 않습니다"),
    ]
    rx = 6.5
    for i, (num, desc) in enumerate(facts):
        ry = 1.85 + i * 1.1
        box(s, rx, ry, 6.23, 0.95, CARD)
        tb(s, num,  rx+0.2, ry+0.15, 1.6, 0.65, 16, bold=True, color=INDIGO)
        tb(s, desc, rx+1.9, ry+0.12, 4.1, 0.7,  11, color=DARK)

    bottom_bar(s, "AI 전문 회사가 아닙니다. 병원 플랫폼 회사입니다.", bg=WARN_L, color=WARN)

    # ── Slide 3: 4단계 진화 ──────────────────────────────────
    s = blank(); bar(s)
    title(s, "AI 활용은 4단계로 진화하고 있다")

    EW = 2.88
    EH = 4.3
    EY = 1.38
    EXS = [0.6 + i * (EW + 0.25) for i in range(4)]

    stages = [
        ("1", "Prompt\nEngineering",
         "채팅창에 질문하기",
         '"이 코드 에러 원인이 뭐야?"',
         "단발성 질의\n창 닫으면 끝"),
        ("2", "Context\nEngineering",
         "CLAUDE.md로 맞춤 세팅",
         "프로젝트 규칙·스타일을\nAI가 매번 기억",
         "한 번 설정\n→ 계속 적용"),
        ("3", "Compound\nEngineering",
         "팀 공유 + 외부 도구 결합",
         "Slack·Notion·DB를\nAI가 직접 실행·조합",
         "개인 → 조직\n레벨 자동화"),
        ("4", "Harness\nEngineering",
         "자율 에이전틱 루프",
         "코드→테스트→에러 수정\nAI가 스스로 반복",
         "24시간 풀 자동화\nCI/CD 결합"),
    ]

    for i, (num, name, sub, desc, bottom) in enumerate(stages):
        is_first = (i == 0)
        cbg = WARN_L if is_first else CARD
        cline = WARN if is_first else None
        box(s, EXS[i], EY, EW, EH, cbg, line_color=cline)
        num_c = WARN if is_first else INDIGO
        tb(s, num, EXS[i]+0.15, EY+0.15, 0.5, 0.5, 22, bold=True, color=num_c)
        name_c = WARN if is_first else INDIGO
        tb(s, name, EXS[i]+0.15, EY+0.6, EW-0.3, 0.85, 15, bold=True, color=name_c)
        tb(s, sub, EXS[i]+0.15, EY+1.5, EW-0.3, 0.45, 11, bold=True, color=DARK)
        tb(s, desc, EXS[i]+0.15, EY+2.05, EW-0.3, 1.0, 11, color=GRAY)
        tb(s, bottom, EXS[i]+0.15, EY+3.2, EW-0.3, 0.9, 10, color=GRAY)
        if i < 3:
            tb(s, "→", EXS[i]+EW+0.02, EY+EH/2-0.25, 0.21, 0.5, 14,
               bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # "YOU ARE HERE" marker
    box(s, EXS[0], EY+EH+0.1, EW, 0.5, WARN)
    tb(s, "대부분이 여기에 있다", EXS[0], EY+EH+0.15, EW, 0.4,
       12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    bottom_bar(s, "채팅 AI에 머물면 임계점을 넘을 수 없다")

    # ── Slide 4: 채팅의 한계 ─────────────────────────────────
    s = blank(); bar(s)
    title(s, "AI를 채팅으로만 쓰면 생기는 일")
    card(s, 0.6, 1.38, 5.93, 4.6, "💬",
         '웹 채팅',
         '"분명히 설명했는데 왜 또 물어봐?"\n\n→ 대화창 닫으면 초기화\n→ 컨텍스트가 길어질수록\n   앞 내용을 잊는다\n→ 매번 처음부터 다시 설명')
    card(s, 6.8, 1.38, 5.93, 4.6, "📱",
         '슬랙',
         '"3개월 전 대화, 찾을 수 있나요?"\n\n→ 빠른 소통의 대가는 휘발\n→ 일주일 전 업무 요청,\n   지금 추적 가능한가요?')
    bottom_bar(s, "채팅은 빠르지만 남지 않는다. 남지 않으면 쌓이지 않는다.")

    # ── Slide 5: Claude Code 비교 ────────────────────────────
    s = blank(); bar(s)
    title(s, "터미널을 켜면 AI의 속이 보인다")
    box(s, 0.6, 1.38, 5.55, 4.95, CARD)
    tb(s, "기존 채팅 방식", 0.9, 1.6, 5.0, 0.55, 15, bold=True, color=GRAY)
    for i, t in enumerate(["대화 → 사라짐", "프롬프트가 중요", "매번 설명 반복", "창 닫으면 초기화"]):
        tb(s, f"✕  {t}", 0.9, 2.38 + i*0.9, 5.0, 0.6, 14, color=GRAY)
    tb(s, "→", 6.35, 3.5, 0.63, 0.75, 30, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    box(s, 7.18, 1.38, 5.55, 4.95, IND_L)
    tb(s, "Claude Code", 7.48, 1.6, 5.0, 0.55, 15, bold=True, color=INDIGO)
    for i, t in enumerate(["대화 → 파일로 저장", "컨텍스트가 중요", "한 번 설정 → 계속 기억", "작업 히스토리가 쌓인다"]):
        tb(s, f"✓  {t}", 7.48, 2.38 + i*0.9, 5.0, 0.6, 14, color=INDIGO)

    # ── Slide 6: Context Engineering ─────────────────────────
    s = blank(); bar(s)
    title(s, "프롬프트 잘 쓰는 시대는 끝났다")
    box(s, 0.6, 1.4, 5.5, 1.7, CARD)
    tb(s, "Prompt Engineering",  0.88, 1.55, 5.0, 0.5, 14, bold=True, color=GRAY)
    tb(s, '"질문을 잘 써야 해"', 0.88, 2.1, 5.0, 0.5, 13, color=GRAY)
    tb(s, "→", 6.35, 1.95, 0.63, 0.8, 28, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    box(s, 7.23, 1.4, 5.5, 1.7, IND_L)
    tb(s, "Context Engineering", 7.5, 1.55, 5.0, 0.5, 14, bold=True, color=INDIGO)
    tb(s, '"정보를 잘 쌓아야 해"', 7.5, 2.1, 5.0, 0.5, 13, color=INDIGO)

    # 3 example cards
    examples = [
        ("📄", "CLAUDE.md",         "프로젝트 규칙·스타일을 파일로 정리 → AI가 매번 읽는다"),
        ("📝", "Notion",            "잘 정리된 페이지를 AI에게 붙여넣기 → 그게 곧 컨텍스트"),
        ("🗂️", "Obsidian / 메모앱", "내 업무 기록을 마크다운으로 관리 → AI가 읽을 수 있다"),
    ]
    for i, (icon, name, desc) in enumerate(examples):
        ey = 3.4 + i * 0.95
        box(s, 0.6, ey, 12.13, 0.82, CARD)
        tb(s, icon, 0.85, ey+0.15, 0.5, 0.5, 18)
        tb(s, name, 1.45, ey+0.15, 2.5, 0.5, 13, bold=True, color=DARK)
        tb(s, desc, 4.2, ey+0.18, 8.3, 0.5, 12, color=GRAY)

    bottom_bar(s, "도구는 바뀐다. Claude Code도 언젠가 사라질 수 있다.\n남는 것은 '내 정보를 AI가 읽을 수 있게 정리하는 습관'이다.",
               y=6.35, h=0.9, size=13)

    # ── Slide 7: 복리 ───────────────────────────────────────
    s = blank(); bar(s)
    title(s, "매일 쌓이면 복리가 된다")
    TW = 3.65
    TH = 4.0
    TY = 1.38
    steps = [
        ("Day 1",  "Claude Code를 켜고,\n아무 업무나 시켜본다"),
        ("Day 7",  "사소한 것부터 모든 업무에\n써본다. 메모도 시키고,\nCLAUDE.md 정리도 시켜본다.\nAI가 나를 어떻게 기록하는지\n지켜본다"),
        ("Day 30", "AI가 더 이상 물어보지\n않는다.\n맥락을 알고 바로 실행한다"),
    ]
    TXS = [0.6 + i * (TW + 0.42) for i in range(3)]
    for i, (day, desc) in enumerate(steps):
        bg = IND_L if i == 2 else CARD
        box(s, TXS[i], TY, TW, TH, bg)
        tb(s, day,  TXS[i]+0.25, TY+0.25, TW-0.5, 0.6, 20, bold=True, color=INDIGO)
        tb(s, desc, TXS[i]+0.25, TY+1.0,  TW-0.5, TH-1.3, 13, color=DARK)
        if i < 2:
            tb(s, "→", TXS[i]+TW+0.08, TY+TH/2-0.3, 0.26, 0.55, 16,
               bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    bottom_bar(s, "임계점을 넘으면, 일하는 방식이 달라진다")

    # ── Slide 8: Phase 1 ────────────────────────────────────
    s = blank(); bar(s)
    title(s, "Phase 1 — 나만의 업무 기록부터")
    CX8 = [0.6, 4.77, 8.94]
    card(s, CX8[0], 1.38, 3.84, 4.6, "📄",
         "내 CLAUDE.md",
         "= 내 업무 설명서\nAI가 나를 기억하는 파일\n나만의 컨텍스트")
    card(s, CX8[1], 1.38, 3.84, 4.6, "📌",
         "매일 commit",
         "= 내 업무 일지\n누가 시키지 않아도\n기록이 쌓인다")
    card(s, CX8[2], 1.38, 3.84, 4.6, "⚡",
         "/wrap",
         "= 자동 정리\n퇴근 전 한 번만 치면\nAI가 알아서 정리")
    bottom_bar(s, "내일 아침, Claude Code를 켜는 것부터 시작입니다", y=6.18, h=0.87, size=17)

    # ── Slide 9a: 공유 지능 (개념) ───────────────────────────
    s = blank(); bar(s)
    title(s, "Phase 2 — 팀이 함께 쌓으면")
    tb(s, "공유 지능", 0.6, 1.22, 12.13, 0.9, 34, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    pairs = [
        ('"지은님이 저번 주에 뭐 하셨죠?"', "AI가 알고 있다"),
        ('"이 업무 누가 담당하고 있죠?"',   "AI가 알고 있다"),
        ('"비슷한 거 예전에 해봤나요?"',    "AI가 찾아준다"),
    ]
    for i, (q, a) in enumerate(pairs):
        cy = 2.32 + i * 1.48
        box(s, 0.6, cy, 8.5, 1.25, CARD)
        tb(s, q, 0.85, cy+0.3, 8.0, 0.65, 14, color=DARK)
        box(s, 9.33, cy, 3.4, 1.25, IND_L)
        tb(s, a, 9.33, cy+0.3, 3.4, 0.65, 14, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    tb(s, "슬랙 대화가 아니라, git 히스토리가 팀의 기억이 됩니다",
       0.6, 6.85, 12.13, 0.5, 13, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 9b: 공유 지능 (구조) ───────────────────────────
    s = blank(); bar(s)
    title(s, "CLAUDE.md를 팀의 오답노트로 쓴다")
    tb(s, "단순히 프롬프트를 잘 쓰는 것을 넘어, 팀의 정보와 AI의 실수를 복리처럼 쌓아가는 핵심 문서",
       0.6, 1.15, 12.13, 0.5, 13, color=GRAY)

    sections = [
        ("1", "우리 팀의 역할",           "목적, 주요 업무 한 줄 정의"),
        ("2", "컨텍스트 경로",            "Notion 링크, 공유 폴더, 성공 사례 모음 경로"),
        ("3", "기본 규칙 · 톤앤매너",     "문서 구조, 어조, 데이터 표기 기준"),
        ("4", "🚨 오답노트",              "AI가 틀릴 때마다 팀원이 추가 → 같은 실수 방지"),
        ("5", "필수 워크플로우",           "시작 전 Plan → 완료 후 검증 → 퇴근 전 /wrap"),
    ]
    for i, (num, sec_name, sec_desc) in enumerate(sections):
        sy = 1.82 + i * 0.92
        is_anti = (i == 3)
        bg_c = WARN_L if is_anti else CARD
        box(s, 0.6, sy, 12.13, 0.8, bg_c, line_color=WARN if is_anti else None)
        tb(s, num, 0.85, sy+0.18, 0.4, 0.44, 14, bold=True, color=WARN if is_anti else INDIGO)
        tb(s, sec_name, 1.4, sy+0.18, 3.5, 0.44, 13, bold=True, color=DARK)
        tb(s, sec_desc, 5.2, sy+0.18, 7.3, 0.44, 12, color=GRAY)

    bottom_bar(s, "팀원이 발견한 AI의 실수가 쌓일수록, AI는 더 똑똑해진다")

    # ── Slide 9c: 공유 지능 (실제 예시) ──────────────────────
    s = blank(); bar(s)
    title(s, "이렇게 쌓인다")

    # code block style
    box(s, 0.6, 1.3, 12.13, 5.0, RGBColor(0x2D, 0x2D, 0x3D))
    tb(s, "🚨 팀 공용 오답노트 (Anti-patterns)", 1.0, 1.55, 11.0, 0.55, 16, bold=True, color=WARN)

    entries = [
        ("[03.15 추가]",
         "경쟁사 분석 시, 2022년 이전 데이터 절대 인용 금지.\n반드시 출처 연도 확인 → 최신 데이터만 반영"),
        ("[03.20 추가]",
         "회의록 요약 시, 사담 제외.\n오직 '결정 사항'과 '다음 할 일'만 추출"),
        ("[04.02 추가]",
         "코드 작성 시, 가짜 API 엔드포인트 생성 금지.\n반드시 사내 공식 API 문서 경로만 사용"),
    ]
    for i, (date, content) in enumerate(entries):
        ey = 2.35 + i * 1.2
        tb(s, date,    1.0, ey,      2.0, 0.4, 12, bold=True, color=GREEN)
        tb(s, content, 3.2, ey,      9.3, 0.9, 12, color=RGBColor(0xCC, 0xCC, 0xDD))

    bottom_bar(s, "한 명이 발견하면, 팀 전체가 같은 실수를 안 한다. 이게 공유 지능이다.")

    # ── Slide 10: 행동 지시 (dark) ───────────────────────────
    s = blank(bg=DK_BG)
    tb(s, "오늘 저녁부터",
       1.2, 0.6,  11.0, 0.55, 20, color=GRAY, align=PP_ALIGN.CENTER)
    tb(s, "딱 이것만",
       1.2, 1.1, 11.0, 0.9, 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    acts = [
        ("①", "Claude Code 설치",                    "내일 아침부터 바로 쓸 수 있도록"),
        ("②", "아침에 PC 켜면 Claude Code도 켠다",   "아무 업무나 시켜본다. 사소한 것부터"),
        ("③", "퇴근 전: /wrap → commit → push",      ""),
    ]
    for i, (num, text, sub) in enumerate(acts):
        cy = 2.35 + i * 0.95
        tb(s, num,  1.2,  cy, 0.6, 0.65, 26, bold=True, color=INDIGO)
        tb(s, text, 2.0,  cy+0.05, 7.0, 0.45, 18, color=WHITE)
        if sub:
            tb(s, sub, 2.0, cy+0.48, 7.0, 0.35, 12, color=GRAY)

    # ③ detail box
    box(s, 1.2, 5.15, 10.93, 1.5, RGBColor(0x25, 0x26, 0x3B))
    details = [
        ("/wrap",   "AI가 오늘 작업을 자동 정리하고, 문서를 갱신하고, 내일 할 일까지 뽑아준다"),
        ("commit",  '오늘 작업의 "저장점"을 찍는다 (문서에 버전 번호 붙이는 것과 같다)'),
        ("push",    "그 저장점을 클라우드에 올린다 (공유 드라이브에 올리는 것과 같다)"),
    ]
    for i, (cmd, desc) in enumerate(details):
        dy = 5.25 + i * 0.42
        tb(s, cmd,  1.5,  dy, 1.2, 0.38, 12, bold=True, color=INDIGO)
        tb(s, "→",  2.75, dy, 0.3, 0.38, 12, color=GRAY)
        tb(s, desc, 3.1,  dy, 8.8, 0.38, 11, color=RGBColor(0xBB, 0xBB, 0xCC))

    tb(s, "이 3개만 매일 반복하면, 컨텍스트의 복리가 시작된다",
       1.2, 6.85, 10.93, 0.45, 15, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    # ── Save ─────────────────────────────────────────────────
    out = os.path.join(os.path.dirname(__file__), "ai_native_intro_v2.pptx")
    prs.save(out)
    print(f"저장 완료: {out}")
    print(f"슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    make_pptx()
