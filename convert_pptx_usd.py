import sys, io, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

RATE = 1450
SRC = 'C:/Users/ash/Downloads/townhall_finance_editable.pptx'
DST = 'C:/Users/ash/Downloads/townhall_finance_usd.pptx'

def krw_to_usd(text):
    def replace(m):
        num = float(m.group(1))
        unit = m.group(2)
        if unit == 'M':
            krw = num * 1_000_000
        elif unit == 'B':
            krw = num * 1_000_000_000
        elif unit == 'K':
            krw = num * 1_000
        usd = krw / RATE
        if usd >= 1_000_000:
            return f'${usd/1_000_000:.1f}M'
        elif usd >= 100_000:
            return f'${round(usd/1000):.0f}K'
        elif usd >= 10_000:
            return f'${usd/1000:.0f}K'
        elif usd >= 1_000:
            return f'${usd/1000:.1f}K'
        else:
            return f'${usd:.0f}'
    return re.sub(r'₩(\d+(?:\.\d+)?)([MBK])', replace, text)

prs = Presentation(SRC)
count = 0

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if '₩' not in para.text:
                continue
            for run in para.runs:
                if '₩' in run.text:
                    old = run.text
                    new = krw_to_usd(old)
                    if old != new:
                        run.text = new
                        count += 1
                        print(f'  Slide {i+1} | {old!r} -> {new!r}')

prs.save(DST)
print(f'\n총 {count}개 변환 완료')
print(f'저장: {DST}')
