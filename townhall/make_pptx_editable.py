"""
EO Studio Town Hall Finance Slides
Editable PowerPoint v3 — English, Zoom-Optimized
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = Path(__file__).parent / "townhall_finance_editable.pptx"

# ── Colors ────────────────────────────────────────────
class C:
    BLACK      = RGBColor(0x0A, 0x0A, 0x0A)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT     = RGBColor(0xE8, 0xFF, 0x47)   # Yellow-green
    CARD       = RGBColor(0x14, 0x14, 0x14)
    GRAY_MID   = RGBColor(0x2A, 0x2A, 0x2A)
    LINE       = RGBColor(0x28, 0x28, 0x28)
    MUTED      = RGBColor(0x66, 0x66, 0x66)
    DIMMED     = RGBColor(0x44, 0x44, 0x44)
    NEON_RED   = RGBColor(0xFF, 0x2D, 0x55)   # Bright neon red
    NEON_GREEN = RGBColor(0x0A, 0xFF, 0x6C)   # Bright neon green

# ── Layout constants (design grid: 1280 × 720 px) ────
W, H, PAD = 1280, 720, 72
CY       = 210   # Y where cards start
CH       = 456   # Card height (fills to footer)
FOOTER_Y = H - 36

def px(p):      return Inches(p * 13.33 / 1280)
def fpt(size):  return Pt(size * 0.75)

# ── Primitives ────────────────────────────────────────
def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = C.BLACK

def txt(slide, text, x, y, w, h, size, color,
        bold=False, align=PP_ALIGN.LEFT, wrap=False):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = fpt(size)
    run.font.bold       = bold
    run.font.color.rgb  = color
    run.font.name       = "Malgun Gothic"
    return box

def rect(slide, x, y, w, h, color, border=None):
    s = slide.shapes.add_shape(1, px(x), px(y), px(w), px(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if border: s.line.color.rgb = border; s.line.width = Pt(0.75)
    else:      s.line.fill.background()
    return s

def hline(slide, x, y, w, color=None):
    r = slide.shapes.add_shape(1, px(x), px(y), px(w), px(1))
    r.fill.solid(); r.fill.fore_color.rgb = color or C.LINE
    r.line.fill.background()

def card_header(slide, cx, cy, cw, label):
    txt(slide, label, cx + 14, cy + 14, cw - 28, 18, 10, C.DIMMED, bold=True)
    hline(slide, cx + 14, cy + 36, cw - 28)

def slide_header(slide, label, page):
    txt(slide, "FINANCE  ·  2026.02", PAD, 14, 400, 18,  9, C.ACCENT, bold=True)
    txt(slide, label,                 PAD, 36, 700, 30, 22, C.ACCENT, bold=True)
    hline(slide, PAD, 72, W - PAD * 2)
    txt(slide, "EO STUDIO TOWNHALL",  PAD,          FOOTER_Y, 300, 18,  8, C.DIMMED)
    txt(slide, f"0{page} / 04", W - PAD - 90, FOOTER_Y,  90, 18,  8, C.DIMMED,
        align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────
# SLIDE 1 — Cash Position
# ─────────────────────────────────────────────────────
def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); slide_header(s, "Cash Position", 1)

    txt(s, "Consolidated Cash Balance", PAD, 82,  600, 20, 12, C.MUTED)
    txt(s, "₩209M",                     PAD, 100, 500, 80, 72, C.WHITE, bold=True)
    txt(s, "▼ ₩154M vs. Last Month",    PAD, 183, 560, 24, 17, C.NEON_RED, bold=True)
    hline(s, PAD, 207, W - PAD * 2)

    CW = (W - PAD * 2 - 16) // 2   # 560
    RX = PAD + CW + 16              # 648
    RH = (CH - 40) // 3             # row height ≈ 138

    # ── Left card: By Entity ──
    rect(s, PAD, CY, CW, CH, C.CARD)
    card_header(s, PAD, CY, CW, "BY ENTITY")

    for i, (name, val) in enumerate([
        ("Korea",   "₩95M"),
        ("US",      "$71K  ·  ₩103M"),
        ("Vietnam", "₩12M"),
    ]):
        ry = CY + 40 + i * RH
        txt(s, name, PAD + 16, ry + 10,  CW - 32, 24, 16, C.MUTED,  bold=True)
        txt(s, val,  PAD + 16, ry + 38,  CW - 32, 54, 40, C.WHITE,  bold=True)
        if i < 2: hline(s, PAD + 16, ry + RH, CW - 32)

    # ── Right card: Net Change ──
    rect(s, RX, CY, CW, CH, C.CARD)
    card_header(s, RX, CY, CW, "NET CHANGE  JAN → FEB")

    ROWY = (CH - 40) // 4  # 4 equal rows filling the card

    for i, (name, val, name_col) in enumerate([
        ("Korea",   "▼ ₩47M",  C.MUTED),
        ("US",      "▼ ₩94M",  C.MUTED),
        ("Vietnam", "▼ ₩13M",  C.MUTED),
        ("Total",   "▼ ₩154M", C.WHITE),
    ]):
        ry     = CY + 40 + i * ROWY
        name_y = ry + (ROWY - 24) // 2
        val_y  = ry + (ROWY - 44) // 2
        txt(s, name, RX + 16,       name_y, 200, 24, 16, name_col,   bold=True)
        txt(s, val,  RX + CW - 220, val_y,  204, 44, 32, C.NEON_RED, bold=True,
            align=PP_ALIGN.RIGHT)
        if i < 3:
            div_color = C.GRAY_MID if i == 2 else None
            hline(s, RX + 16, ry + ROWY, CW - 32, div_color)

# ─────────────────────────────────────────────────────
# SLIDE 2 — Revenue
# ─────────────────────────────────────────────────────
def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); slide_header(s, "Revenue", 2)

    txt(s, "Jan–Feb Cumulative Revenue",                         PAD, 82,  600, 20, 12, C.MUTED)
    txt(s, "₩432M",                                              PAD, 100, 500, 80, 72, C.WHITE, bold=True)
    txt(s, "Global 94% Contribution  ·  YoY -8% (Same BU Base)", PAD, 183, 800, 22, 13, C.MUTED)
    hline(s, PAD, 207, W - PAD * 2)

    CW = (W - PAD * 2 - 16) // 2
    RX = PAD + CW + 16

    # ── Left card: Monthly Revenue table ──
    rect(s, PAD, CY, CW, CH, C.CARD)
    card_header(s, PAD, CY, CW, "MONTHLY REVENUE")

    # Table column headers
    TY = CY + 46
    txt(s, "BU",    PAD + 16,  TY, 150, 18, 11, C.DIMMED, bold=True)
    txt(s, "Jan",   PAD + 178, TY, 100, 18, 11, C.DIMMED, bold=True, align=PP_ALIGN.RIGHT)
    txt(s, "Feb",   PAD + 296, TY, 100, 18, 11, C.DIMMED, bold=True, align=PP_ALIGN.RIGHT)
    txt(s, "Total", PAD + 414, TY, 114, 18, 11, C.DIMMED, bold=True, align=PP_ALIGN.RIGHT)
    hline(s, PAD + 16, TY + 24, CW - 32)

    ROWY = (CH - 76) // 4   # 4 rows (3 data + 1 total)
    for i, (name, v1, v2, vtot) in enumerate([
        ("KR Content", "₩1.6M",   "₩13.1M", "₩14.7M"),
        ("KR Planet",  "₩6.5M",   "₩5.0M",  "₩11.5M"),
        ("GL Global",  "₩320.7M", "₩85.2M", "₩405.9M"),
    ]):
        ry = CY + 76 + i * ROWY + 6
        txt(s, name,  PAD + 16,  ry, 160, 26, 16, C.WHITE, bold=True)
        txt(s, v1,    PAD + 178, ry, 100, 26, 16, C.MUTED, bold=True, align=PP_ALIGN.RIGHT)
        txt(s, v2,    PAD + 296, ry, 100, 26, 16, C.WHITE, bold=True, align=PP_ALIGN.RIGHT)
        txt(s, vtot,  PAD + 414, ry, 114, 26, 16, C.ACCENT, bold=True, align=PP_ALIGN.RIGHT)
        if i < 2: hline(s, PAD + 16, CY + 76 + (i + 1) * ROWY, CW - 32)

    sty = CY + 76 + 3 * ROWY
    hline(s, PAD + 16, sty, CW - 32, C.GRAY_MID)
    txt(s, "Total",    PAD + 16,  sty + 8, 160, 28, 18, C.WHITE,  bold=True)
    txt(s, "₩328.9M",  PAD + 178, sty + 8, 100, 28, 18, C.WHITE,  bold=True, align=PP_ALIGN.RIGHT)
    txt(s, "₩103.3M",  PAD + 296, sty + 8, 100, 28, 18, C.WHITE,  bold=True, align=PP_ALIGN.RIGHT)
    txt(s, "₩432.2M",  PAD + 414, sty + 8, 114, 28, 18, C.ACCENT, bold=True, align=PP_ALIGN.RIGHT)

    # ── Right card: YoY ──
    rect(s, RX, CY, CW, CH, C.CARD)
    card_header(s, RX, CY, CW, "YOY  2025 vs 2026 (Same BU)")

    YOWY = (CH - 40) // 4
    for i, (name, ref25, pct, color) in enumerate([
        ("GL Global",  "2025:  \u20a9319M", "+27%", C.NEON_GREEN),
        ("KR Content", "2025:  \u20a963M",  "-77%", C.NEON_RED),
        ("KR Planet",  "2025:  \u20a987M",  "-87%", C.NEON_RED),
    ]):
        ry = CY + 40 + i * YOWY
        txt(s, name,  RX + 16,       ry + 12, 260, 24, 16, C.WHITE,  bold=True)
        txt(s, ref25, RX + 16,       ry + 38, 260, 20, 12, C.DIMMED)
        txt(s, pct,   RX + CW - 240, ry + 4,  224, 58, 44, color,    bold=True,
            align=PP_ALIGN.RIGHT)
        if i < 2: hline(s, RX + 16, ry + YOWY, CW - 32)

    oty = CY + 40 + 3 * YOWY
    hline(s, RX + 16, oty, CW - 32, C.GRAY_MID)
    txt(s, "Overall (Same BU)", RX + 16,       oty + 10, 300, 30, 20, C.WHITE,    bold=True)
    txt(s, "-8%",               RX + CW - 240, oty + 2,  224, 60, 46, C.NEON_RED, bold=True,
        align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────
# SLIDE 3 — BU Goal Achievement
# ─────────────────────────────────────────────────────
def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); slide_header(s, "Revenue Goal Achievement", 3)

    txt(s, "Annual Goal Achievement (Confirmed Deals)",               PAD, 82,  700, 20, 12, C.MUTED)
    txt(s, "17.8%",                                                   PAD, 100, 500, 80, 72, C.WHITE, bold=True)
    txt(s, "Annual Target: ₩12B  |  Confirmed: ₩2.1B  |  Remaining: ₩9.9B",
        PAD, 183, 900, 22, 13, C.MUTED)
    hline(s, PAD, 207, W - PAD * 2)

    CGAP = 14
    CW   = (W - PAD * 2 - CGAP * 2) // 3  # ≈ 370

    for i, c in enumerate([
        {"title": "KR",       "target": "₩1.3B", "conf": "₩330M",  "gap": "-₩970M",  "pct": 0.254, "ps": "25.4%"},
        {"title": "PLANNING", "target": "₩2.2B", "conf": "₩850M",  "gap": "-₩1.35B", "pct": 0.388, "ps": "38.8%"},
        {"title": "GLOBAL",   "target": "₩8.5B", "conf": "₩960M",  "gap": "-₩7.5B",  "pct": 0.112, "ps": "11.2%"},
    ]):
        cx = PAD + i * (CW + CGAP)
        rect(s, cx, CY, CW, CH, C.CARD)
        txt(s, c["title"], cx + 14, CY + 14, CW - 28, 22, 14, C.ACCENT, bold=True)
        hline(s, cx + 14, CY + 40, CW - 28)

        # Target
        txt(s, "TARGET",    cx + 14, CY + 52,  CW - 28, 18, 11, C.DIMMED, bold=True)
        txt(s, c["target"], cx + 14, CY + 70,  CW - 28, 64, 50, C.WHITE,  bold=True)

        # Confirmed
        txt(s, "CONFIRMED", cx + 14, CY + 144, CW - 28, 18, 11, C.DIMMED, bold=True)
        txt(s, c["conf"],   cx + 14, CY + 162, CW - 28, 56, 42, C.WHITE,  bold=True)

        # Progress bar
        BARY = CY + 228
        BARW = CW - 28
        rect(s, cx + 14, BARY, BARW,                 10, C.GRAY_MID)
        rect(s, cx + 14, BARY, int(BARW * c["pct"]), 10, C.ACCENT)

        # Achievement %
        txt(s, c["ps"], cx + 14, BARY + 18, CW - 28, 64, 50, C.ACCENT, bold=True)

        # Gap
        txt(s, f"Gap  {c['gap']}", cx + 14, BARY + 90, CW - 28, 22, 14, C.DIMMED, bold=True)

        # KR Content footnote
        if i == 0:
            txt(s, "* Includes \u20a9300M KTO Production",
                cx + 14, BARY + 116, CW - 28, 18, 9, C.DIMMED)

# ─────────────────────────────────────────────────────
# SLIDE 4 — AI Native
# ─────────────────────────────────────────────────────
def slide4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); slide_header(s, "AI Native", 4)

    # Giant headline — ~40% of slide height
    txt(s, "The Era of",     PAD + 40,  76, 1000, 136, 102, C.WHITE,  bold=True)
    txt(s, "10x Leverage.",  PAD + 40, 208, 1000, 136, 102, C.ACCENT, bold=True)

    # Body
    txt(s,
        "AI is not here to replace you. It's a lever to amplify your capabilities.",
        PAD + 40, 362, 960, 36, 20, C.MUTED, bold=True, wrap=True)
    txt(s,
        "Build your own 24/7 AI team to clone your workflow. That is AI Native.",
        PAD + 40, 406, 960, 46, 20, C.WHITE, bold=True, wrap=True)

    # CTA button
    cta = s.shapes.add_shape(1, px(PAD + 40), px(472), px(880), px(60))
    cta.fill.background()
    cta.line.color.rgb = C.ACCENT
    cta.line.width     = Pt(1)
    txt(s, "\u2192   Install Claude Code & Build Your AI Team", PAD + 72, 488, 840, 34, 20, C.ACCENT, bold=True)

# ─────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide1(prs); print("Slide 1: Cash Position done")
    slide2(prs); print("Slide 2: Revenue done")
    slide3(prs); print("Slide 3: BU Goal Achievement done")
    slide4(prs); print("Slide 4: AI Native done")

    prs.save(str(OUTPUT))
    print(f"\nDone → {OUTPUT}")

if __name__ == "__main__":
    main()
