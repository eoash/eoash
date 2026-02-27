"""
PNG 4장을 16:9 PowerPoint 파일로 만듭니다.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SLIDES_DIR = Path(__file__).parent
OUTPUT     = SLIDES_DIR / "townhall_finance.pptx"

PNG_FILES = [
    SLIDES_DIR / "slide1.png",
    SLIDES_DIR / "slide2.png",
    SLIDES_DIR / "slide3.png",
    SLIDES_DIR / "slide4.png",
]

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃

    for i, png in enumerate(PNG_FILES, 1):
        if not png.exists():
            print(f"파일 없음: {png}")
            continue

        slide = prs.slides.add_slide(blank_layout)

        # 배경 검정
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0A)

        # PNG를 슬라이드 전체 크기로 삽입
        slide.shapes.add_picture(
            str(png),
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"  Slide {i} 추가: {png.name}")

    prs.save(str(OUTPUT))
    print(f"\n완료 → {OUTPUT}")

if __name__ == "__main__":
    main()
